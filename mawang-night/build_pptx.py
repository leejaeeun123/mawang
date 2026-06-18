#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
RED HOUR — 마왕족발 브랜드 경험 프로젝트
편집 가능한 네이티브 .pptx 생성기 (python-pptx)

  pip install python-pptx
  python build_pptx.py     ->  RED_HOUR.pptx

모든 텍스트는 실제 편집 가능한 run 으로, 견적은 네이티브 표로 생성된다.
사진/로고는 자리표(점선 박스)로 표시 → PowerPoint 에서 이미지로 교체.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image as PILImage
import os, tempfile

# ---- palette ----
INK      = RGBColor(0x0A, 0x0A, 0x0B)
INK2     = RGBColor(0x15, 0x15, 0x18)
PANEL    = RGBColor(0x1B, 0x1B, 0x1E)
REDPANEL = RGBColor(0x1F, 0x10, 0x12)
BLOOD    = RGBColor(0xC4, 0x14, 0x1B)
BLOODDP  = RGBColor(0x7A, 0x0A, 0x0E)
BLOODSF  = RGBColor(0xE2, 0x3B, 0x3F)
BONE     = RGBColor(0xF4, 0xF0, 0xE8)
ASH      = RGBColor(0x9C, 0x9A, 0xA0)
ASHD     = RGBColor(0x6C, 0x6A, 0x70)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0x33, 0x24, 0x26)

FONT = "Pretendard"
EMU_IN = 914400
PW, PH = 11.69, 8.27   # A4 landscape inches

prs = Presentation()
prs.slide_width  = Inches(PW)
prs.slide_height = Inches(PH)
BLANK = prs.slide_layouts[6]


# ---------- low-level helpers ----------
def _no_line(shape):
    shape.line.fill.background()

def rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp

def dashed(slide, x, y, w, h):
    sp = rect(slide, x, y, w, h, REDPANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, line_w=1.0)
    ln = sp.line._get_or_add_ln()
    d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(d)
    return sp

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
_TMPD = tempfile.mkdtemp()
def place_image(slide, x, y, w, h, fname):
    """center-crop image to the box aspect ratio and insert at exact size."""
    im = PILImage.open(os.path.join(ASSETS, fname)).convert("RGB")
    tgt = (w / h); iw, ih = im.size; cur = iw / ih
    if cur > tgt:
        nw = int(ih * tgt); l = (iw - nw) // 2; im = im.crop((l, 0, l + nw, ih))
    else:
        nh = int(iw / tgt); t = (ih - nh) // 2; im = im.crop((0, t, iw, t + nh))
    out = os.path.join(_TMPD, fname.replace('/', '_'))
    im.save(out, quality=88)
    pic = slide.shapes.add_picture(out, Inches(x), Inches(y), Inches(w), Inches(h))
    pic.line.color.rgb = LINE; pic.line.width = Pt(0.75)
    return pic

def _set_run(r, text, size, color, bold=False, italic=False, spacing=None):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    # ensure east-asian font too
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT}); rPr.append(ea)
    cs = rPr.makeelement(qn('a:cs'), {'typeface': FONT}); rPr.append(cs)
    if spacing is not None:
        rPr.set('spc', str(int(spacing * 100)))
    return r

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, 0)
    tf.paragraphs[0].alignment = align
    return tb, tf

def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=2, space_before=0, line=1.0, bullet=False, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    try:
        p.line_spacing = line
    except Exception:
        pass
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    if bullet:
        # diamond-ish bullet using a red dash char
        b = pPr.makeelement(qn('a:buChar'), {'char': '▪'});
        f = pPr.makeelement(qn('a:buFont'), {'typeface': FONT})
        clr = pPr.makeelement(qn('a:buClr'), {})
        srgb = clr.makeelement(qn('a:srgbClr'), {'val': 'C4141B'}); clr.append(srgb)
        pPr.append(clr); pPr.append(f); pPr.append(b)
        pPr.set('indent', '-137160'); pPr.set('marL', '137160')
    else:
        n = pPr.makeelement(qn('a:buNone'), {}); pPr.append(n)
    return p


# ---------- slide scaffolding ----------
ML = 1.0           # content left margin (in)
CW = PW - ML - 0.6 # content width

def new_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, PW, PH, INK)                    # background
    rect(s, 0, 0, 0.34, PH, BLOOD)                # red edge bar
    return s

def header(s, kicker, pno):
    # wordmark
    tb, tf = textbox(s, ML, 0.42, 4.5, 0.4, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, first=True)
    _set_run(p.add_run(), "마왕", 14, BONE, bold=True)
    _set_run(p.add_run(), "족발", 14, BLOOD, bold=True)
    _set_run(p.add_run(), "   RED HOUR", 8, ASHD, bold=True, spacing=2)
    # kicker + page no (right)
    tb2, tf2 = textbox(s, PW-5.1, 0.42, 4.5, 0.4, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    p2 = para(tf2, first=True, align=PP_ALIGN.RIGHT)
    _set_run(p2.add_run(), kicker + "    ", 11, BLOODSF, bold=True, spacing=2)
    _set_run(p2.add_run(), pno, 11, ASHD, bold=True, spacing=1)

def footer(s, pno):
    rect(s, ML, PH-0.62, CW, 0.012, LINE)
    tb, tf = textbox(s, ML, PH-0.55, CW, 0.3)
    p = para(tf, first=True)
    _set_run(p.add_run(), "RED HOUR · 마왕족발 브랜드 경험 프로젝트", 7.5, ASHD, spacing=1)
    tb2, tf2 = textbox(s, PW-2.0, PH-0.55, 1.4, 0.3, align=PP_ALIGN.RIGHT)
    p2 = para(tf2, first=True, align=PP_ALIGN.RIGHT)
    _set_run(p2.add_run(), pno, 9, BLOOD, bold=True, spacing=1)

def secno(s, text, y=1.15):
    tb, tf = textbox(s, ML, y, CW, 0.3)
    p = para(tf, first=True)
    _set_run(p.add_run(), text, 11, BLOOD, bold=True, spacing=3)

def title(s, parts, y=1.45, size=26):
    """parts: list of (text, color)"""
    tb, tf = textbox(s, ML, y, CW, 1.0)
    p = para(tf, first=True, line=1.08)
    for text, color in parts:
        _set_run(p.add_run(), text, size, color, bold=True)
    return tb

def lead(s, text, y, size=12, w=None, align=PP_ALIGN.LEFT):
    tb, tf = textbox(s, ML, y, w or CW, 1.1)
    p = para(tf, first=True, line=1.4, align=align)
    _set_run(p.add_run(), text, size, ASH)
    return tb

def card(s, x, y, w, h, ctitle, bullets, red=False, bignum=None):
    bg = REDPANEL if red else PANEL
    sp = rect(s, x, y, w, h, bg, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              line=(BLOOD if red else LINE), line_w=0.75)
    tb, tf = textbox(s, x+0.18, y+0.14, w-0.36, h-0.26)
    p = para(tf, first=True, space_after=5)
    _set_run(p.add_run(), "◆ ", 9, BLOOD, bold=True)
    _set_run(p.add_run(), ctitle, 12, BONE, bold=True)
    for b in bullets:
        bp = para(tf, bullet=True, line=1.32, space_after=2)
        # allow (text, bold) tuples inside via simple **marker
        _set_run(bp.add_run(), b, 10.5, ASH)
    return sp

def cards_row(s, items, top, height, gap=0.18):
    n = len(items)
    w = (CW - gap*(n-1)) / n
    for i, it in enumerate(items):
        x = ML + i*(w+gap)
        card(s, x, top, w, height, it[0], it[1], red=it[2] if len(it) > 2 else False)


# ============================================================ SLIDES

# ---- P01 COVER ----
def cover():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, PW, PH, INK)
    rect(s, 0, 0, 0.34, PH, BLOOD)
    # logo placeholder top-right
    dashed(s, PW-2.4, 0.5, 1.8, 0.75)
    tbl, tfl = textbox(s, PW-2.4, 0.5, 1.8, 0.75, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    pl = para(tfl, first=True, align=PP_ALIGN.CENTER, line=1.2)
    _set_run(pl.add_run(), "마왕족발 로고\n삽입 영역", 8.5, ASHD)
    # texts
    tb, tf = textbox(s, ML, 2.4, 8.5, 4.0)
    p = para(tf, first=True, space_after=8)
    _set_run(p.add_run(), "BRAND EXPERIENCE PROJECT", 15, BLOOD, bold=True, spacing=5)
    p2 = para(tf, line=0.95, space_after=2)
    _set_run(p2.add_run(), "RED ", 66, BONE, bold=True)
    _set_run(p2.add_run(), "HOUR", 66, BLOOD, bold=True)
    p3 = para(tf, space_after=10)
    _set_run(p3.add_run(), "레드 아워 ", 20, ASH, bold=True, spacing=5)
    _set_run(p3.add_run(), "(가제)", 12, ASHD)
    p4 = para(tf, space_after=8)
    _set_run(p4.add_run(), "마왕족발 ", 22, BONE, bold=True)
    _set_run(p4.add_run(), "브랜드 경험", 22, BLOOD, bold=True)
    _set_run(p4.add_run(), " 프로젝트", 22, BONE, bold=True)
    p5 = para(tf, line=1.4, space_before=6)
    _set_run(p5.add_run(), "정해진 한 시간, 도시의 밤을 마왕의 레드로 물들인다.\n마왕족발의 세계관을 ‘경험’으로 전환하는 브랜드 나이트.", 12, ASH)
    p6 = para(tf, space_before=14)
    _set_run(p6.add_run(), "마왕족발 × RED HOUR   |   ", 11, ASH, spacing=1)
    _set_run(p6.add_run(), "2026", 11, BLOOD, bold=True, spacing=1)
    footer(s, "01")

# ---- generic title+cards slide ----
def slide_cards(kicker, pno, sec, title_parts, lead_text, cards, t_size=28,
                cards_top=2.55, cards_h=3.2, lead_y=None):
    s = new_slide()
    header(s, kicker, pno)
    secno(s, sec)
    title(s, title_parts, size=t_size)
    top = cards_top
    if lead_text:
        ly = lead_y if lead_y else 2.25
        lead(s, lead_text, ly)
    cards_row(s, cards, top, cards_h)
    footer(s, pno)
    return s

# ---- P02 ----
def p02():
    s = new_slide(); header(s, "PROJECT OVERVIEW", "02")
    secno(s, "01 — 프로젝트 개요")
    title(s, [("왜 ", BONE), ("지금", BLOOD), (", 레드 아워인가", BONE)], size=26)
    cards_row(s, [
        ("강한 캐릭터·세계관", ["마왕족발은 이미 뚜렷한 캐릭터와 세계관 보유", "먹거리를 넘어 ‘경험·문화’ 브랜드로 확장할 시점"]),
        ("팬덤 결집 필요", ["온라인 팬덤을 오프라인에서 결집할 접점 부재", "찐팬·인플루언서와의 직접 경험 설계 필요"]),
        ("콘텐츠 = 자산", ["현장 장면이 곧 바이럴 콘텐츠·UGC가 되는 시대", "한 번의 행사가 지속 가능한 자산으로 축적"]),
    ], 2.2, 1.9)
    secno(s, "프로젝트 목표", y=4.35)
    goals = [("경험 각인","브랜드를 ‘경험’으로 각인"),("팬덤 결집","오프라인 팬덤 결집"),
             ("콘텐츠 확보","고품질 UGC 대량 확보"),("제품 테스트","신메뉴·콜라보 주류 검증"),
             ("사업 확장","차기 팝업·굿즈 레퍼런스")]
    n=5; gap=0.16; w=(CW-gap*(n-1))/n
    for i,(t,b) in enumerate(goals):
        x=ML+i*(w+gap)
        card(s, x, 4.75, w, 1.55, t, [b], red=True)
    footer(s, "02")

# ---- P03 ----
def p03():
    s = new_slide(); header(s, "RE-FRAMING", "03")
    secno(s, "02 — 기획 방향 재정립")
    title(s, [("공실 팝업이 아니라, ", BONE), ("브랜드 나이트", BLOOD), ("로", BONE)], size=24)
    card(s, ML, 2.4, CW/2-0.1, 3.0, "기존 공실 팝업 (한계)", [
        "별도 공간 대관 + 집기·장비 렌탈 필요",
        "운영 스태프·바텐더·DJ 전부 신규 섭외",
        "장시간 운영 → 비용↑ 운영 피로도↑",
        "결국 단순 판매에 머무를 위험",
    ])
    card(s, ML+CW/2+0.1, 2.4, CW/2-0.1, 3.0, "RED HOUR (브랜드 나이트)", [
        "라운지·클럽의 인프라·인력 그대로 활용",
        "하루에 집중 → 화제성·참여·확산 폭발",
        "판매가 아닌 경험·콘텐츠 중심 설계",
        "비용 대비 브랜드 임팩트 극대화",
    ], red=True)
    lead(s, "전환의 핵심 — “오래 여는 팝업”에서 “강렬한 하룻밤의 브랜드 경험”으로.", 5.7, align=PP_ALIGN.CENTER)
    footer(s, "03")

# ---- P04 ----
def p04():
    s = new_slide(); header(s, "CONCEPT", "04")
    secno(s, "03 — 행사 컨셉")
    tb, tf = textbox(s, ML, 1.6, CW, 1.5)
    p = para(tf, first=True, line=1.2)
    _set_run(p.add_run(), "“가장 마왕스럽게 ", 30, BONE, bold=True)
    _set_run(p.add_run(), "경험하는 시간", 30, BLOOD, bold=True)
    _set_run(p.add_run(), ", RED HOUR.”", 30, BONE, bold=True)
    lead(s, "도시의 밤, 정해진 한 시간에 마왕족발의 세계관을 폭발적으로 경험시킨다. 레드(RED)는 마왕의 시그니처 컬러이자 강렬함·중독·열기의 상징. 그 시간 안에서 음식·공간·사람·콘텐츠가 하나의 장면으로 완성된다.", 3.1)
    kws = [("INTENSE","강렬하고 중독적인 한 시간"),("RED","마왕의 시그니처 레드 무드"),("ONE NIGHT","하루에 집중된 브랜드 경험")]
    n=3; gap=0.2; w=(CW-gap*(n-1))/n
    for i,(t,d) in enumerate(kws):
        x=ML+i*(w+gap)
        rect(s, x, 4.7, w, 1.7, REDPANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=BLOOD, line_w=0.75)
        tbk, tfk = textbox(s, x, 4.95, w, 1.25, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        pk = para(tfk, first=True, align=PP_ALIGN.CENTER, space_after=6)
        _set_run(pk.add_run(), t, 24, BLOODSF, bold=True, spacing=1)
        pk2 = para(tfk, align=PP_ALIGN.CENTER)
        _set_run(pk2.add_run(), d, 10, ASH)
    footer(s, "04")

# ---- P05 ----
def p05():
    s = new_slide(); header(s, "SPACE", "05")
    secno(s, "04 — 공간 방향성")
    title(s, [("라운지·클럽 ", BONE), ("협업형", BLOOD), (" 공간", BONE)], size=26)
    lead(s, "자체 공간을 새로 세팅하는 대신, 이미 갖춰진 라운지·클럽과 협업해 운영 리스크는 낮추고 분위기는 극대화한다.", 2.2)
    cards_row(s, [
        ("인프라 활용", ["집기·서빙 인력·운영 시스템 그대로 연계", "바텐더·DJ·음향 장비 기구축"]),
        ("분위기 자산", ["조명·사운드·동선 등 공간의 무드 활용", "레드 연출로 마왕 세계관과 합치"]),
        ("브랜드는 경험에 집중", ["운영 부담 덜고 콘텐츠·경험 설계 집중", "검토 공간은 부록 참조 (어덜트온리/로우키)"], True),
    ], 3.0, 1.9)
    place_image(s, ML, 5.05, CW/2-0.1, 1.35, "venue_lowkey_a.jpg")
    place_image(s, ML+CW/2+0.1, 5.05, CW/2-0.1, 1.35, "venue_adult_a.jpg")
    footer(s, "05")

def ph_label(s, x, y, w, h, text):
    tb, tf = textbox(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    _set_run(p.add_run(), "🖼  " + text, 9.5, ASHD)

# ---- P06 ----
def p06():
    s = new_slide(); header(s, "DRESS CODE", "06")
    secno(s, "05 — 드레스 코드")
    title(s, [("DRESS CODE : ", BONE), ("BLACK OR RED", BLOOD)], size=26)
    lead(s, "입장객 전원에게 블랙 또는 레드 드레스코드를 적용한다. 통일된 컬러는 현장을 하나의 장면으로 묶고, 사진·영상의 완성도와 소속감을 끌어올린다.", 2.2)
    kws = [("BLACK","마왕의 어두운 세계관", BONE),("RED","강렬한 시그니처 레드", BLOODSF),("BLACK + RED","가장 마왕스러운 조합", WHITE)]
    n=3; gap=0.2; w=(CW-gap*(n-1))/n
    for i,(t,d,c) in enumerate(kws):
        x=ML+i*(w+gap)
        rect(s, x, 3.1, w, 1.55, REDPANEL if i!=0 else INK2, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=BLOOD if i==1 else LINE)
        tbk, tfk = textbox(s, x, 3.3, w, 1.15, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        pk = para(tfk, first=True, align=PP_ALIGN.CENTER, space_after=5)
        _set_run(pk.add_run(), t, 22, c, bold=True, spacing=1)
        pk2 = para(tfk, align=PP_ALIGN.CENTER)
        _set_run(pk2.add_run(), d, 10, ASH)
    card(s, ML, 4.95, CW/2-0.1, 1.4, "현장 연출", ["레드 조명 + 통일 컬러 → 포토제닉한 공간","컬러별 포토존·입장 연출 운영"])
    card(s, ML+CW/2+0.1, 4.95, CW/2-0.1, 1.4, "참여 베네핏", ["드레스코드 충족 시 입장 혜택·기프트","SNS 인증 → 자연 확산 트리거"], red=True)
    footer(s, "06")

# ---- P07 ----
def p07():
    s = new_slide(); header(s, "STAFF", "07")
    secno(s, "06 — 스태프 컨셉")
    title(s, [("스태프가 곧 ", BONE), ("마왕 세계관", BLOOD)], size=26)
    lead(s, "현장 스태프를 단순 운영 인력이 아니라 마왕 세계관의 캐릭터로 연출한다. 살아 움직이는 캐릭터는 그 자체로 강력한 브랜드 경험이자 포토 콘텐츠가 된다.", 2.2)
    cards_row(s, [
        ("캐릭터 스타일링", ["데빌·좀비 무드의 분장·의상","블랙·레드 톤으로 통일"]),
        ("소품·디테일", ["마왕 시그니처 소품·액세서리","존(VIP/바/포토존)별 역할 분장"]),
        ("경험 효과", ["입장 순간부터 세계관 몰입","스태프와의 포토·인증 콘텐츠 유발"], True),
    ], 3.0, 1.9)
    dashed(s, ML, 5.15, CW, 1.2)
    ph_label(s, ML, 5.15, CW, 1.2, "스태프 분장·스타일링 레퍼런스 · 삽입 영역")
    footer(s, "07")

# ---- P08 timeline ----
def p08():
    s = new_slide(); header(s, "SCHEDULE", "08")
    secno(s, "07 — 행사 일정")
    title(s, [("SATURDAY ONLY, ", BONE), ("세 개의 장면", BLOOD)], size=24)
    lead(s, "토요일 단 하루, 하나의 밤을 세 개의 장면으로 설계한다.", 2.25)
    slots = [
        ("18:00–19:00","VIP EXPERIENCE",["초청 30인 내외","신메뉴 선체험 · 콜라보 주류","스타일링 테이블 · 포토타임","VIP 전용 기프트"]),
        ("19:00–20:00","RED HOUR EXPERIENCE",["일반 고객 입장","신메뉴 체험 · 콜라보 주류","굿즈 증정 (키캡 키링)","VIP 기대감 → 일반 확장"]),
        ("21:00–01:00","MAWANG NIGHT PARTY",["DJ 공연 · 게스트","전문 바텐더 · 믹스밤 쇼잉","프리드링크 · 상품권 추첨","자유 파티 / 네트워킹"]),
    ]
    n=3; gap=0.2; w=(CW-gap*(n-1))/n; top=2.9; h=2.7
    for i,(tm,nm,items) in enumerate(slots):
        x=ML+i*(w+gap)
        rect(s, x, top, w, h, PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE)
        rect(s, x, top, w, 0.78, BLOOD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tbh, tfh = textbox(s, x+0.16, top+0.08, w-0.3, 0.66, anchor=MSO_ANCHOR.MIDDLE)
        ph = para(tfh, first=True, space_after=1)
        _set_run(ph.add_run(), tm, 17, WHITE, bold=True)
        ph2 = para(tfh)
        _set_run(ph2.add_run(), nm, 11, WHITE, bold=True)
        tbb, tfb = textbox(s, x+0.16, top+0.92, w-0.32, h-1.0)
        for j,it in enumerate(items):
            pb = para(tfb, first=(j==0), bullet=True, line=1.3, space_after=3)
            _set_run(pb.add_run(), it, 10, ASH)
    lead(s, "※ 20:00–21:00 인터미션 — 공간 전환 및 파티 세팅 / 현장 전체 전문 촬영으로 콘텐츠 아카이빙.", 5.85, size=10.5)
    footer(s, "08")

# ---- P09 ----
def p09():
    s = new_slide(); header(s, "18:00 — VIP", "09")
    secno(s, "07-1 — VIP EXPERIENCE (18:00–19:00)")
    title(s, [("초청 30인, ", BONE), ("기대감의 진원지", BLOOD)], size=24)
    cards_row(s, [
        ("대상 (30인 내외)", ["마왕족발 먹방 인플루언서","점주 추천 ‘찐팬’ 고객","브랜드 관계자","인스타 신청자 중 선정"]),
        ("프로그램", ["신메뉴 선체험","마왕 콜라보 주류 시음","스타일링된 테이블 연출","자유 숏폼 촬영 / 포토타임","브랜드 소개(짧고 강하게)"]),
        ("VIP GIFT", ["MAWANG KEYCAP KEYRING 증정","VIP 전용 추가 기프트","찐팬 대상 스페셜 굿즈"], True),
    ], 2.5, 3.0)
    lead(s, "※ VIP 세션 전체를 전문 촬영 → 본편 파티의 기대감을 만드는 선행 콘텐츠로 활용.", 5.8, size=11)
    footer(s, "09")

# ---- P10 ----
def p10():
    s = new_slide(); header(s, "19:00 — EXPERIENCE", "10")
    secno(s, "07-2 — RED HOUR EXPERIENCE (19:00–20:00)")
    title(s, [("기대감을 ", BONE), ("모두의 경험", BLOOD), ("으로", BONE)], size=24)
    card(s, ML, 2.5, CW/2-0.1, 2.5, "입장 / 운영", [
        "인스타 광고·사전 신청을 통해 일반 고객 입장","드레스코드(BLACK OR RED) 적용","VIP 세션의 무드를 그대로 이어받아 확장"])
    card(s, ML+CW/2+0.1, 2.5, CW/2-0.1, 2.5, "경험 콘텐츠", [
        "신메뉴 체험 · 콜라보 주류 시음","굿즈 증정 — MAWANG KEYCAP KEYRING","포토존·스태프 연출로 인증 콘텐츠 유발"], red=True)
    dashed(s, ML, 5.25, CW, 1.1)
    ph_label(s, ML, 5.25, CW, 1.1, "현장 경험 컷 (신메뉴·주류·포토존) · 삽입 영역")
    footer(s, "10")

# ---- P11 ----
def p11():
    s = new_slide(); header(s, "21:00 — PARTY", "11")
    secno(s, "07-3 — MAWANG NIGHT PARTY (21:00–01:00)")
    title(s, [("밤의 본편, ", BONE), ("바이럴이 터지는 시간", BLOOD)], size=24)
    card(s, ML, 2.5, CW/2-0.1, 2.5, "프로그램", [
        "DJ 공연","래퍼 게스트 공연 (가능 시)","전문 바텐더 운영","마왕 레드 사케 믹스밤 쇼잉","자유 파티 / 네트워킹"])
    card(s, ML+CW/2+0.1, 2.5, CW/2-0.1, 2.5, "이벤트 (바이럴 트리거)", [
        "입장객 전원 프리 드링크 1잔","콜라보 주류 증정 이벤트","마왕족발 상품권 추첨","현장 인증 시 추가 혜택"], red=True)
    dashed(s, ML, 5.25, CW, 1.1)
    ph_label(s, ML, 5.25, CW, 1.1, "파티 무드 / DJ·바텐더 컷 · 삽입 영역")
    footer(s, "11")

# ---- P12 ----
def p12():
    s = new_slide(); header(s, "FOOD", "12")
    secno(s, "08 — 푸드 경험")
    title(s, [("파티에 맞춘 ", BONE), ("시그니처 푸드", BLOOD)], size=26)
    lead(s, "기존 족발을 파티·핑거푸드 형태로 재해석해, 마시고 즐기는 밤에 어울리는 두 가지 시그니처를 제안한다.", 2.2)
    card(s, ML, 2.95, CW/2-0.1, 2.3, "① DEVIL COLD SKEWER", [
        "차갑게 즐기는 마왕 시그니처 콜드 꼬치","강렬한 비주얼 + 마왕 특유의 매운맛","한입 사이즈로 술과 페어링·포토제닉"], red=True)
    card(s, ML+CW/2+0.1, 2.95, CW/2-0.1, 2.3, "② FIRE CUP", [
        "한 손에 즐기는 컵 푸드","서서 즐기는 파티 친화형 포맷","레드 비주얼로 인증 콘텐츠 유발"])
    dashed(s, ML, 5.45, CW/2-0.1, 0.9); ph_label(s, ML, 5.45, CW/2-0.1, 0.9, "DEVIL COLD SKEWER 사진")
    dashed(s, ML+CW/2+0.1, 5.45, CW/2-0.1, 0.9); ph_label(s, ML+CW/2+0.1, 5.45, CW/2-0.1, 0.9, "FIRE CUP 사진")
    footer(s, "12")

# ---- P13 ----
def p13():
    s = new_slide(); header(s, "GOODS", "13")
    secno(s, "09 — 굿즈 전략")
    title(s, [("갖고 싶은 ", BONE), ("마왕 굿즈", BLOOD)], size=26)
    card(s, ML, 2.5, CW/2-0.1, 2.6, "MAWANG KEYCAP KEYRING", [
        "마왕 캐릭터 키캡을 끼운 DIY 키캡 키링","직접 조립하는 참여형 굿즈","VIP·일반 세션 공통 증정 → 참여·확산 트리거","수량·구성·단가는 견적 부록 참조"], red=True)
    card(s, ML+CW/2+0.1, 2.5, CW/2-0.1, 1.6, "기타 굿즈", [
        "스티커팩 — 마왕 캐릭터·세계관","거울 등 데일리 소품","블랙·레드 무드로 비주얼 톤 통일"])
    dashed(s, ML+CW/2+0.1, 4.25, CW/2-0.1, 0.85)
    ph_label(s, ML+CW/2+0.1, 4.25, CW/2-0.1, 0.85, "굿즈 시안 이미지 삽입 영역")
    footer(s, "13")

# ---- P14 ----
def p14():
    s = new_slide(); header(s, "CONTENT", "14")
    secno(s, "10 — 콘텐츠 전략")
    tb, tf = textbox(s, ML, 1.55, CW, 0.9)
    p = para(tf, first=True)
    _set_run(p.add_run(), "“이벤트가 아니라, ", 28, BONE, bold=True)
    _set_run(p.add_run(), "콘텐츠다.”", 28, BLOOD, bold=True)
    lead(s, "미션형 인스타 이벤트는 최소화하고, 공간·분위기·장면 자체가 콘텐츠가 되도록 설계한다.", 2.5)
    cards_row(s, [
        ("사전", ["먹방 인플루언서 사전 섭외","숏폼 제작 가이드(필수 컷)","티저 콘텐츠 선공개"]),
        ("현장", ["전문 촬영팀 운영","행사 전체 기록·아카이빙","포토존·스태프 연출로 UGC 유발"]),
        ("사후", ["종료 후 바이럴 광고 집행","릴스·숏폼 재편집 배포","고품질 UGC·아카이브 축적"], True),
    ], 3.2, 2.6)
    footer(s, "14")

# ---- P15 ----
def p15():
    s = new_slide(); header(s, "AFTER", "15")
    secno(s, "11 — 행사 이후 활용")
    title(s, [("하룻밤을 ", BONE), ("브랜드 자산", BLOOD), ("으로", BONE)], size=26)
    lead(s, "RED HOUR는 하루로 끝나지 않는다. 현장에서 만든 모든 장면을 지속 가능한 브랜드 자산으로 재활용한다.", 2.2)
    cards_row(s, [
        ("마케팅 자산", ["릴스·숏폼 광고 집행","브랜드 SNS 정기 업로드"]),
        ("사업 자산", ["점주 대상 홍보 자료","신메뉴·주류 론칭 콘텐츠"]),
        ("확장 자산", ["차기 팝업·콜라보 레퍼런스","굿즈·IP 사업 테스트 베드"], True),
    ], 3.0, 2.4)
    footer(s, "15")

# ---- P16 ----
def p16():
    s = new_slide(); header(s, "IMPACT", "16")
    secno(s, "12 — 기대 효과")
    title(s, [("하루의 임팩트, ", BONE), ("세 방향의 효과", BLOOD)], size=26)
    cards_row(s, [
        ("브랜드", ["이미지 리프레시","경험 브랜드로 각인","찐팬 로열티 강화"]),
        ("마케팅", ["고품질 UGC·바이럴","MZ 접점 확대","콘텐츠 아카이브 구축"]),
        ("사업", ["신메뉴·콜라보 주류 테스트","굿즈·IP 사업 가능성 검증","차기 사업 레퍼런스 확보"], True),
    ], 2.5, 2.8)
    footer(s, "16")

# ---- P17 ----
def p17():
    s = new_slide(); header(s, "KEY MESSAGE", "17")
    rect(s, ML, 2.0, 1.7, 0.42, REDPANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=BLOOD)
    tbt, tft = textbox(s, ML, 2.0, 1.7, 0.42, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    pt = para(tft, first=True, align=PP_ALIGN.CENTER)
    _set_run(pt.add_run(), "KEY MESSAGE", 10, BLOOD, bold=True, spacing=2)
    tb, tf = textbox(s, ML, 2.9, CW, 2.4)
    p = para(tf, first=True, line=1.3)
    _set_run(p.add_run(), "“RED HOUR는 팝업이 아니라,\n마왕족발을 ", 28, BONE, bold=True)
    _set_run(p.add_run(), "‘경험’으로 전환", 28, BLOOD, bold=True)
    _set_run(p.add_run(), "하는\n브랜드 나이트다.”", 28, BONE, bold=True)
    lead(s, "가장 마왕스러운 한 시간이, 브랜드의 다음을 만든다.", 5.4)
    footer(s, "17")

# ---- table helper ----
def style_cell(cell, text, size, color, bold=False, fill=None, align=PP_ALIGN.LEFT):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = INK2
    cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    _set_run(p.add_run(), text, size, color, bold=bold)

# ---- A1 venue compare ----
def a1():
    s = new_slide(); header(s, "APPENDIX · VENUE", "A1")
    secno(s, "APPENDIX")
    title(s, [("공간 대여 검토 ", BONE), ("2곳", BLOOD)], size=24)
    rows = [
        ("", "① 어덜트 온리 (ADULT ONLY)", "② 로우키 (LOWKEY)"),
        ("위치", "가로수길", "압구정로데오"),
        ("인스타", "@adultonlybar", "@lowkeyseoul"),
        ("특징", "레드·고급, 인테리어가 이미 마왕과 어울리는 편", "공간이 넓어 마왕 스타일로 꾸미기 수월"),
        ("대관 / 비용", "평일 600만 / 주말 800만 (VAT별도, 조정 가능)", "단순 대관 시간당 30만 · 팝업+파티 시 조정"),
        ("주말 대관", "가능 (조건 협의)", "까다롭지 않음"),
        ("브랜드 콜라보", "가능", "가능"),
        ("인력·DJ", "서버·바텐더 시간당 16,000원 / DJ 15~20만", "현장 직원·DJ 바로 연결"),
    ]
    rn=len(rows); cn=3
    tbl = s.shapes.add_table(rn, cn, Inches(ML), Inches(2.3), Inches(CW), Inches(3.6)).table
    tbl.columns[0].width = Inches(1.7)
    tbl.columns[1].width = Inches((CW-1.7)/2)
    tbl.columns[2].width = Inches((CW-1.7)/2)
    for ci in range(cn):
        style_cell(tbl.cell(0,ci), rows[0][ci], 11, WHITE, bold=True, fill=BLOOD)
    for ri in range(1, rn):
        style_cell(tbl.cell(ri,0), rows[ri][0], 10, BLOODSF, bold=True, fill=PANEL)
        style_cell(tbl.cell(ri,1), rows[ri][1], 9.5, ASH, fill=INK2)
        style_cell(tbl.cell(ri,2), rows[ri][2], 9.5, ASH, fill=INK2)
    lead(s, "※ 비용·조건은 협의 진행 중인 참고용 수치이며, 규모·시간·범위에 따라 조정될 수 있음.", 6.15, size=10.5)
    footer(s, "A1")

# ---- venue detail (A2/A3) ----
def venue_detail(pno, kicker, head_parts, ig, blocks, chat_title, chats, ph_text, ph_image=None):
    s = new_slide(); header(s, kicker, pno)
    secno(s, "APPENDIX — 공간 " + ("①" if pno=="A2" else "②"))
    title(s, head_parts, size=22)
    # left blocks
    lx = ML; lw = CW/2-0.15; y = 2.35
    for (bt, items, red) in blocks:
        h = 0.55 + 0.255*len(items)
        card(s, lx, y, lw, h, bt, items, red=red)
        y += h + 0.12
    # right: chat
    rx = ML+CW/2+0.15; rw = CW/2-0.15
    tb, tf = textbox(s, rx, 2.35, rw, 0.3)
    pc = para(tf, first=True)
    _set_run(pc.add_run(), "◆ " + chat_title, 12, BONE, bold=True)
    cy = 2.75
    for who, side, text in chats:
        h = 0.52 + 0.2*max(1, (len(text)//26))
        bx = rx if side=="q" else rx+rw*0.12
        bw = rw*0.88
        fill = PANEL if side=="q" else REDPANEL
        ln = LINE if side=="q" else BLOOD
        rect(s, bx, cy, bw, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=ln, line_w=0.5)
        tbb, tfb = textbox(s, bx+0.12, cy+0.07, bw-0.24, h-0.12)
        pw = para(tfb, first=True, space_after=1)
        _set_run(pw.add_run(), who, 8, BLOODSF, bold=True, spacing=1)
        pb = para(tfb, line=1.25)
        _set_run(pb.add_run(), text, 9.5, BONE if side=="a" else ASH)
        cy += h + 0.1
    ph_h = 6.35 - cy
    if ph_h > 0.4:
        if ph_image:
            place_image(s, rx, cy+0.05, rw, ph_h-0.05, ph_image)
        else:
            dashed(s, rx, cy+0.02, rw, ph_h); ph_label(s, rx, cy+0.02, rw, ph_h, ph_text)
    footer(s, pno)

def a2():
    venue_detail("A2", "APPENDIX · ① ADULT ONLY",
        [("어덜트 온리 ", BONE), ("(가로수길)", BLOOD)], "@adultonlybar",
        [
          ("공간 / 운영", ["레드·고급 무드, 인테리어가 마왕과 잘 어울림","지하 1·2층 복층 · 단독 출입구 (가드 배치 고려)","공간 변경 가능 · 브랜드 콜라보 가능","신메뉴 시식·협업 주류, 주방·바 활용 가능"], False),
          ("비용 / 인력 (참고용)", ["대관: 평일 600만 / 주말·공휴일전일 800만 (VAT별도)","규모·시간·방문객·협업 범위에 따라 조정","서버 & 바텐더 시간당 16,000원","서울 DJ(1:30 기준) 150,000~200,000원"], True),
          ("활용 아이디어", ["프리드링크·족발체험·키캡/상품권 증정·DJ 파티","SNS 대대적 광고 → 마왕 찐팬 파티로 운영"], False),
        ],
        "사장님 소통 요약",
        [("MAWANG","q","주말 단독 대관·브랜드 콜라보 가능할까요? 비용은요?"),
         ("ADULT ONLY","a","운영은 초대제 중심. 단독 대관 시 초대자만 입장하는 방식이 좋아요."),
         ("ADULT ONLY","a","평일 600 / 주말 800(VAT별도) 기준, 범위 따라 조정 가능합니다."),
         ("ADULT ONLY","a","주방·바 활용, 협업 주류 OK. 서버·바텐더 시간당 16,000원, DJ 15~20만.")],
        "어덜트 온리 공간 사진", ph_image="venue_adult_b.jpg")

def a3():
    venue_detail("A3", "APPENDIX · ② LOWKEY",
        [("로우키 ", BONE), ("(압구정로데오)", BLOOD)], "@lowkeyseoul",
        [
          ("공간 / 운영", ["공간이 넓어 마왕 스타일 연출 수월","현장 집기·직원·DJ 바로 연결","주말 대관 무난 · 브랜드 콜라보 가능","조명 밝기·색상 조절 가능"], False),
          ("비용 / 조건 (참고용)", ["단순 대관 시간당 30만","팝업+파티 시 보장 매출에 따라 조정 가능","푸드트럭 주차공간 있으나 1층 매장 옆 → 협의"], True),
          ("시즌 / 연출 팁", ["성수기 지금~가을, 관광객 한겨울 빼고 꾸준 (장마철 회피)","밤 파티 시 주류 서버·DJ 필요, 레드 사케 믹스밤 고려"], False),
        ],
        "대표(주현성) 소통 요약",
        [("MAWANG","q","주말 팝업+파티 대관 가능? 조명·푸드트럭 세팅도 되나요?"),
         ("주현성 대표","a","넓어서 마왕 스타일로 꾸미기 좋아요. 조명 조절 가능, 직원·DJ 바로 붙습니다."),
         ("주현성 대표","a","단순 대관 시간당 30만, 팝업+파티면 보장 매출 맞춰 조정 가능."),
         ("주현성 대표","a","푸드트럭은 1층 매장 옆이라 협의 필요, 시즌은 장마철만 피하면 무난.")],
        "로우키 공간 사진", ph_image="venue_lowkey_b.jpg")

# ---- A4 quote 뉴리디파인 ----
def a4():
    s = new_slide(); header(s, "APPENDIX · QUOTE ①", "A4")
    secno(s, "APPENDIX — 굿즈 견적 ① 뉴리디파인")
    title(s, [("MAWANG KEYCAP KEYRING ", BONE), ("견적", BLOOD)], size=20)
    head = ["구성","품목","수량","단가","금액"]
    data = [
        ("구성 A (3구)","3구 키링 (투명/블랙/화이트)","1,000","2,900","2,900,000"),
        ("","풀컬러 키캡 (OEM, 체리)","3,000","1,100","3,300,000"),
        ("","개별 조립·포장 (3구)","1,000","600","600,000"),
        ("구성 B (2구)","2구 키링","1,000","2,400","2,400,000"),
        ("","풀컬러 키캡 (OEM, 체리)","2,000","1,100","2,200,000"),
        ("","개별 조립·포장 (2구)","1,000","400","400,000"),
        ("구성 C (1구)","1구 키링 (투명)","1,000","1,200","1,200,000"),
        ("","풀컬러 키캡 (OEM, 체리)","1,000","1,100","1,100,000"),
        ("","개별 조립·포장 (1구)","1,000","200","200,000"),
    ]
    foots = [("구성 A 합계 (VAT 포함)","6,800,000"),("구성 B 합계 (VAT 포함)","5,000,000"),("구성 C 합계 (VAT 포함)","2,500,000")]
    rn = 1 + len(data) + len(foots); cn = 5
    tbl = s.shapes.add_table(rn, cn, Inches(ML), Inches(2.1), Inches(CW), Inches(3.7)).table
    widths = [1.6, 4.0, 1.2, 1.3, 1.99]
    for i,wd in enumerate(widths): tbl.columns[i].width = Inches(wd)
    aligns = [PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.RIGHT,PP_ALIGN.RIGHT,PP_ALIGN.RIGHT]
    for ci in range(cn):
        style_cell(tbl.cell(0,ci), head[ci], 11, WHITE, bold=True, fill=BLOOD, align=aligns[ci])
    for ri,row in enumerate(data, start=1):
        for ci,val in enumerate(row):
            col = BONE if ci==0 else ASH
            style_cell(tbl.cell(ri,ci), val, 9.5, col, bold=(ci==0), align=aligns[ci])
    fr = 1+len(data)
    for k,(lbl,amt) in enumerate(foots):
        ri = fr+k
        tbl.cell(ri,0).merge(tbl.cell(ri,3))
        style_cell(tbl.cell(ri,0), lbl, 11, BONE, bold=True, fill=REDPANEL, align=PP_ALIGN.LEFT)
        style_cell(tbl.cell(ri,4), amt, 11, BONE, bold=True, fill=REDPANEL, align=PP_ALIGN.RIGHT)
    lead(s, "비고 · 풀컬러 키캡 디자인 5종까지 할인 · 대량 제작 보통 2~3주(매우 많으면 소폭 증가) · 개별 조립·포장 제외 시 해당 금액만 차감.", 6.0, size=10)
    footer(s, "A4")

# ---- A5 quote 브랜다즐 ----
def a5():
    s = new_slide(); header(s, "APPENDIX · QUOTE ②", "A5")
    secno(s, "APPENDIX — 굿즈 견적 ② 브랜다즐 (BRANDAZZLE)")
    title(s, [("MAWANG KEYCAP KEYRING ", BONE), ("견적 (비교)", BLOOD)], size=20)
    head = ["구분","품목","수량","단가","공급가액"]
    data = [
        ("미조립형","1구 보드(LED)+개별 PVC케이스","1,000","1,700","1,700,000"),
        ("","키캡 (4종×250)/미조립","1,000","1,000","1,000,000"),
        ("","2구 보드(LED)+개별 PVC케이스","1,000","2,300","2,300,000"),
        ("","키캡 (4종×500)/미조립","2,000","900","1,800,000"),
        ("","3구 보드(LED)+개별 PVC케이스","1,000","3,000","3,000,000"),
        ("","키캡 (4종×750)/미조립","3,000","850","2,550,000"),
        ("완성형","1구 키캡키링(LED)/OPP (4종×250)","1,000","2,400","2,400,000"),
        ("","2구 키캡키링(LED)/OPP (4종×500)","1,000","4,000","4,000,000"),
        ("","3구 키캡키링(LED)/OPP (4종×750)","1,000","5,500","5,500,000"),
    ]
    foots = [("미조립형 총금액 (VAT 포함)","1구 2,970,000 · 2구 4,510,000 · 3구 6,105,000"),
             ("완성형 총금액 (VAT 포함)","1구 2,640,000 · 2구 4,400,000 · 3구 6,050,000")]
    rn = 1 + len(data) + len(foots); cn = 5
    tbl = s.shapes.add_table(rn, cn, Inches(ML), Inches(2.0), Inches(CW), Inches(3.5)).table
    widths = [1.3, 4.6, 1.0, 1.0, 2.19]
    for i,wd in enumerate(widths): tbl.columns[i].width = Inches(wd)
    aligns = [PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.RIGHT,PP_ALIGN.RIGHT,PP_ALIGN.RIGHT]
    for ci in range(cn):
        style_cell(tbl.cell(0,ci), head[ci], 10.5, WHITE, bold=True, fill=BLOOD, align=aligns[ci])
    for ri,row in enumerate(data, start=1):
        for ci,val in enumerate(row):
            style_cell(tbl.cell(ri,ci), val, 9, BONE if ci==0 else ASH, bold=(ci==0), align=aligns[ci])
    # group label merges
    tbl.cell(1,0).merge(tbl.cell(6,0)); style_cell(tbl.cell(1,0), "미조립형\n(보드+키캡)", 9.5, BONE, bold=True, fill=PANEL)
    tbl.cell(7,0).merge(tbl.cell(9,0)); style_cell(tbl.cell(7,0), "완성형\n(키링 LED·OPP)", 9.5, BONE, bold=True, fill=PANEL)
    fr = 1+len(data)
    for k,(lbl,amt) in enumerate(foots):
        ri = fr+k
        tbl.cell(ri,0).merge(tbl.cell(ri,1))
        tbl.cell(ri,2).merge(tbl.cell(ri,4))
        style_cell(tbl.cell(ri,0), lbl, 10, BONE, bold=True, fill=REDPANEL, align=PP_ALIGN.LEFT)
        style_cell(tbl.cell(ri,2), amt, 10, BONE, bold=True, fill=REDPANEL, align=PP_ALIGN.RIGHT)
    lead(s, "비고 · 공급가액 VAT별도(합계 VAT포함) · OPP+made in china 부착 · 견적 영업일 기준 2주 유효 · 결제 발주 후 100%(협의가능) · 발주·디자인파일 영업일 4~5주 · 공급: 주식회사 에이블러(brandazzle_cs@brandazzle.kr)", 5.75, size=9.5)
    footer(s, "A5")


# ---------- build ----------
cover()
p02(); p03(); p04(); p05(); p06(); p07(); p08(); p09(); p10()
p11(); p12(); p13(); p14(); p15(); p16(); p17()
a1(); a2(); a3(); a4(); a5()

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "RED_HOUR.pptx")
prs.save(out)
print("PPTX written:", out, "·", len(prs.slides._sldIdLst), "slides")
